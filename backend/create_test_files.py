import os
import docx
import fitz
import pptx
from PIL import Image, ImageDraw


def make_test_txt():
    content = "This is a simple text file.\n\nIt contains a second paragraph for reading tests.\n\nThird paragraph has more text to verify word counting."
    with open("test.txt", "w", encoding="utf-8") as file_object:
        file_object.write(content)


def make_test_md():
    content = "# Markdown Test Document\n\nThis is paragraph one from markdown.\n\n## Subheading Section\n\nThis is paragraph two explaining markdown content."
    with open("test.md", "w", encoding="utf-8") as file_object:
        file_object.write(content)


def make_test_docx():
    document = docx.Document()
    document.add_heading("DOCX Test Document", level=0)
    document.add_paragraph("This is the first paragraph of the DOCX test file.")
    document.add_paragraph("This is the second paragraph with additional text for validation.")
    document.save("test.docx")


def make_test_pdf():
    document = fitz.open()
    page = document.new_page()
    text = "Adaptive Reader PDF Test Document\n\nThis is the first paragraph on page one of the PDF.\n\nThis is the second paragraph providing reading content."
    page.insert_text((50, 72), text, fontsize=12)
    document.save("test.pdf")
    document.close()


def make_test_pptx():
    presentation = pptx.Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Presentation Slide 1"
    body.text = "This is slide content paragraph one.\nThis is slide content paragraph two."
    presentation.save("test.pptx")


def make_test_image():
    image = Image.new("RGB", (400, 100), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)
    draw.text((10, 30), "Hello OCR Text", fill=(0, 0, 0))
    image.save("test.jpg")


if __name__ == "__main__":
    make_test_txt()
    make_test_md()
    make_test_docx()
    make_test_pdf()
    make_test_pptx()
    make_test_image()
    print("Test files created successfully.")
