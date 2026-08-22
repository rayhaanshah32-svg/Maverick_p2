import os

from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

# ---------------------------------------------------------
# TESSERACT CONFIGURATION
# ---------------------------------------------------------

# Windows Tesseract installation path
TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

if os.path.exists(TESSERACT_PATH):
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH


# ---------------------------------------------------------
# HANDWRITING OCR
# ---------------------------------------------------------

# Import TrOCR module.
# If the handwriting module is unavailable, the normal
# OCR functions will still work.

try:
    from handwriting_ocr import extract_handwriting

    HANDWRITING_OCR_AVAILABLE = True

except ImportError:
    HANDWRITING_OCR_AVAILABLE = False


# ---------------------------------------------------------
# PDF TEXT EXTRACTION
# ---------------------------------------------------------

def extract_pdf(file_path):
    """
    Extract typed/selectable text from a PDF.
    Returns a tuple of (plain_text, page_count).
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(
            f"PDF not found: {file_path}"
        )

    reader = PdfReader(file_path)
    page_count = len(reader.pages)

    pages = []

    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            pages.append(page_text.strip())

    plain_text = "\n\n".join(pages)

    return plain_text, page_count


def extract_docx(file_path):
    """
    Extract text from a Word DOCX file.
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(
            f"DOCX not found: {file_path}"
        )

    document = Document(file_path)

    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    return "\n\n".join(paragraphs)


# ---------------------------------------------------------
# PPTX TEXT EXTRACTION
# ---------------------------------------------------------

def extract_ppt(file_path):
    """
    Legacy alias — calls extract_pptx and returns only the text.
    """
    plain_text, _slide_count = extract_pptx(file_path)
    return plain_text


def extract_pptx(file_path):
    """
    Extract text from PowerPoint slides.
    Returns a tuple of (plain_text, slide_count).
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(
            f"PowerPoint file not found: {file_path}"
        )

    presentation = Presentation(file_path)
    slide_count = len(presentation.slides)

    slides = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text and shape.text.strip()
        ]

        if texts:
            slides.append(
                f"--- SLIDE {slide_number} ---\n" + "\n".join(texts)
            )

    plain_text = "\n\n".join(slides)

    return plain_text, slide_count


def extract_image_ocr(file_path):
    """
    Extract printed text from an image using Tesseract OCR.
    Alias for extract_image used by main.py.
    """
    return extract_image(file_path)


def read_text_file(file_path):
    """
    Read a plain text or Markdown file.
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(
            f"Text file not found: {file_path}"
        )

    with open(file_path, "r", encoding="utf-8", errors="replace") as file:
        return file.read()


# ---------------------------------------------------------
# PRINTED IMAGE OCR
# ---------------------------------------------------------

def extract_image(file_path):
    """
    Extract printed text from JPG/JPEG/PNG
    using Tesseract OCR.
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(
            f"Image not found: {file_path}"
        )

    image = Image.open(file_path)

    image = image.convert("RGB")

    text = pytesseract.image_to_string(
        image
    )

    return text.strip()


# ---------------------------------------------------------
# IMAGE OCR WITH HANDWRITING FALLBACK
# ---------------------------------------------------------

def extract_image_with_handwriting(file_path):
    """
    Try Tesseract first.

    If Tesseract does not find useful text,
    use TrOCR handwriting recognition.
    """

    # -----------------------------------------
    # STEP 1: TESSERACT
    # -----------------------------------------

    printed_text = extract_image(
        file_path
    )

    if printed_text.strip():

        return printed_text.strip()

    # -----------------------------------------
    # STEP 2: TROCR
    # -----------------------------------------

    if HANDWRITING_OCR_AVAILABLE:

        try:

            handwriting_text = extract_handwriting(
                file_path
            )

            if handwriting_text:

                return handwriting_text.strip()

        except Exception as error:

            print(
                "Handwriting OCR failed:",
                error
            )

    return ""


# ---------------------------------------------------------
# PDF OCR WITH HANDWRITING FALLBACK
# ---------------------------------------------------------

def extract_pdf_with_handwriting(file_path):
    """
    First attempt normal PDF text extraction.

    If the PDF has no selectable text,
    use TrOCR on the rendered PDF pages.
    """

    # -----------------------------------------
    # STEP 1: NORMAL PDF EXTRACTION
    # -----------------------------------------

    text = extract_pdf(file_path)

    if text.strip():

        return text.strip()

    # -----------------------------------------
    # STEP 2: HANDWRITING OCR
    # -----------------------------------------

    if HANDWRITING_OCR_AVAILABLE:

        try:

            handwriting_text = extract_handwriting(
                file_path
            )

            if handwriting_text:

                return handwriting_text.strip()

        except Exception as error:

            print(
                "Handwriting PDF OCR failed:",
                error
            )

    return ""


# ---------------------------------------------------------
# MAIN EXTRACTION FUNCTION
# ---------------------------------------------------------

def extract_text(file_path):
    """
    Automatically determine the file type
    and extract text.

    Supported:

    PDF
    PPTX
    JPG
    JPEG
    PNG
    """

    if not os.path.exists(file_path):

        raise FileNotFoundError(
            f"File not found: {file_path}"
        )

    extension = os.path.splitext(
        file_path
    )[1].lower()


    # -----------------------------------------
    # PDF
    # -----------------------------------------

    if extension == ".pdf":

        return extract_pdf_with_handwriting(
            file_path
        )


    # -----------------------------------------
    # POWERPOINT
    # -----------------------------------------

    elif extension == ".pptx":

        return extract_ppt(
            file_path
        )


    # -----------------------------------------
    # IMAGES
    # -----------------------------------------

    elif extension in [
        ".jpg",
        ".jpeg",
        ".png"
    ]:

        return extract_image_with_handwriting(
            file_path
        )


    # -----------------------------------------
    # UNSUPPORTED FILE
    # -----------------------------------------

    else:

        raise ValueError(
            "Unsupported file type: "
            f"{extension}\n"
            "Supported formats: "
            "PDF, PPTX, JPG, JPEG, PNG"
        )


# ---------------------------------------------------------
# SAVE TEXT TO FILE
# ---------------------------------------------------------

def extract_and_save(
    input_file,
    output_file
):
    """
    Extract text and save it into a TXT file.
    """

    print("\n" + "=" * 70)

    print(
        "INPUT FILE:",
        input_file
    )

    print("=" * 70)

    try:

        text = extract_text(
            input_file
        )

        if text:

            print(
                "\nEXTRACTED TEXT:\n"
            )

            print(text)

            with open(
                output_file,
                "w",
                encoding="utf-8"
            ) as file:

                file.write(text)

            print(
                "\nTEXT SAVED SUCCESSFULLY"
            )

            print(
                "OUTPUT FILE:",
                output_file
            )

        else:

            print(
                "\nNo text was detected."
            )

    except Exception as error:

        print(
            "\nERROR:",
            error
        )


# ---------------------------------------------------------
# DIRECT TESTING
# ---------------------------------------------------------

if __name__ == "__main__":

    print(
        "Adaptive Reader - Text Extraction Module"
    )

    print(
        "\nSupported formats:"
    )

    print(
        "PDF | PPTX | JPG | JPEG | PNG"
    )